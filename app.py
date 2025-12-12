import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
import io
import json
import re
import unicodedata
import time

# --- CONFIGURATION ---
st.set_page_config(page_title="Gemini Pro-Translator", page_icon="💎")

# --- CUSTOM CSS FOR DYNAMIC DRAG-AND-DROP ---
st.markdown("""
<style>
    /* 1. Target the main Drag and Drop Section */
    [data-testid='stFileUploader'] {
        width: 100%;
    }
    
    /* 2. The actual dropzone box - Default State */
    [data-testid='stFileUploader'] section {
        background-color: #f8f9fa; /* Light grey background */
        border: 2px dashed #dfe1e5;
        border-radius: 12px;
        padding: 30px;
        transition: all 0.3s ease-in-out;
    }

    /* 3. HARD TEXT COLOR FIX: Force all text inside to be dark grey 
       This overrides the "white on hover" issue for spans, divs, and small text */
    [data-testid='stFileUploader'] section > div,
    [data-testid='stFileUploader'] section span,
    [data-testid='stFileUploader'] section small {
        color: #31333F !important;
    }

    /* 4. THE DYNAMIC EFFECT: When dragging/hovering */
    [data-testid='stFileUploader'] section:hover {
        background-color: #e3f2fd; /* Light Blue Background */
        border: 2px dashed #2196f3; /* Bright Blue Border */
        transform: scale(1.02); /* Slight pop/zoom effect */
        box-shadow: 0 4px 15px rgba(33, 150, 243, 0.3); /* Blue glow shadow */
    }

    /* 5. Button Styling - Force White Text */
        [data-testid='stFileUploader'] button {
            color: white !important;
            background-color: #4CAF50; /* Add a color (like green or dark grey) so white text shows up */
            border: none;
        }

    /* 6. Pulsing animation for the border */
    @keyframes border-pulse {
        0% { border-color: #2196f3; }
        50% { border-color: #64b5f6; }
        100% { border-color: #2196f3; }
    }
    
    [data-testid='stFileUploader'] section:hover {
        animation: border-pulse 2s infinite;
    }
</style>
""", unsafe_allow_html=True)

st.title("💎 Gemini PPT Translator")

# --- API KEY ---
api_key = None
try:
    if "GOOGLE_API_KEY" in st.secrets:
        api_key = st.secrets["GOOGLE_API_KEY"]
except:
    pass 

with st.sidebar:
    if not api_key:
        api_key = st.text_input("Google API Key", type="password")
    
    st.divider()
    st.write("### ⚙️ Scaling Strategy")
    scaling_mode = st.radio("Mode", ["Smart Adaptive Curve", "Manual Fixed"])
    
    manual_scale = 0.75
    if scaling_mode == "Manual Fixed":
        manual_scale = st.slider("Fixed Scale Factor", 0.3, 1.5, 0.75, 0.05)
    else:
        st.caption("Uses a non-linear curve to prevent text from getting too small on dense slides.")

# --- HELPERS ---

def get_text_visual_width(text):
    if not text: return 0
    width = 0
    for char in text:
        if unicodedata.east_asian_width(char) in ['W', 'F', 'A']:
            width += 1.6 
        else:
            width += 1.0 
    return width

def calculate_slide_global_ratio(translation_map):
    total_orig_width = 0
    total_trans_width = 0
    
    for original, translated in translation_map.items():
        total_orig_width += get_text_visual_width(original)
        total_trans_width += get_text_visual_width(translated)
        
    if total_orig_width == 0 or total_trans_width == 0:
        return 1.0
        
    raw_ratio = total_orig_width / total_trans_width
    
    # Adaptive Curve Logic
    if raw_ratio < 0.6:
        adjusted_ratio = raw_ratio * 1.35
    elif raw_ratio < 0.75:
        adjusted_ratio = raw_ratio * 1.25
    else:
        adjusted_ratio = raw_ratio * 1.1
    
    if adjusted_ratio > 1.1: adjusted_ratio = 1.1
    if adjusted_ratio < 0.65: adjusted_ratio = 0.65
    
    return adjusted_ratio

def apply_scaling(source_run, target_run, target_font_name, scale_factor):
    src_font = source_run.font
    dst_font = target_run.font

    if target_font_name == "Montserrat" and src_font.bold:
        dst_font.name = "Montserrat SemiBold"
        dst_font.bold = False 
    else:
        dst_font.name = target_font_name
        dst_font.bold = src_font.bold

    dst_font.italic = src_font.italic
    dst_font.underline = src_font.underline

    base_size = src_font.size if src_font.size else Pt(18)
    
    raw_new_size = base_size.pt * scale_factor
    final_size = int(round(raw_new_size))
    
    if final_size < 4: final_size = 4
        
    dst_font.size = Pt(final_size)

    try:
        if src_font.color.type == MSO_COLOR_TYPE.RGB:
            dst_font.color.rgb = src_font.color.rgb
        elif src_font.color.type == MSO_COLOR_TYPE.SCHEME:
            dst_font.color.theme_color = src_font.color.theme_color
    except: pass 

def get_font_size_pt(text_frame):
    try:
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            run = text_frame.paragraphs[0].runs[0]
            if run.font.size:
                return run.font.size.pt
    except: pass
    return 12.0

def auto_widen_columns(shape):
    if not shape.has_table: return
    table = shape.table
    
    for col_idx, column in enumerate(table.columns):
        max_word_width_emu = 0
        current_col_width = column.width
        
        for row in table.rows:
            if col_idx >= len(row.cells): continue
            cell = row.cells[col_idx]
            txt = cell.text_frame.text.strip()
            if not txt: continue
            
            font_size = get_font_size_pt(cell.text_frame)
            words = txt.split()
            if not words: continue
            
            for word in words:
                word_visual_width = get_text_visual_width(word)
                est_word_width = (word_visual_width * 0.6) * font_size * 12700
                if est_word_width > max_word_width_emu:
                    max_word_width_emu = est_word_width
        
        if max_word_width_emu > current_col_width:
            try:
                expansion_cap = 3.0 if col_idx == 0 else 1.5
                limit = current_col_width * expansion_cap
                new_width = int(max_word_width_emu * 1.25)
                column.width = min(new_width, int(limit))
            except: pass

def fix_text_overflow_centered(shape, translated_text):
    if not shape or not translated_text: return
    try:
        words = translated_text.split()
        if not words: return
        
        font_size = get_font_size_pt(shape.text_frame)
        longest_word = max(words, key=len)
        visual_units = get_text_visual_width(longest_word)
        
        estimated_width_pts = (visual_units * 0.6) * font_size
        required_width_emu = estimated_width_pts * 12700 * 1.2
        
        current_width = shape.width
        
        if required_width_emu > current_width:
            limit = current_width * 2.5
            new_width = min(int(required_width_emu), int(limit))
            
            width_diff = new_width - current_width
            new_left = shape.left - (width_diff // 2)
            
            shape.left = int(new_left)
            shape.width = int(new_width)
    except: pass

# --- STANDARD FUNCTIONS ---
def iter_text_frames(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_frames(shape.shapes)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        yield (shape, cell.text_frame)
        elif shape.has_text_frame:
            yield (shape, shape.text_frame)

def enforce_layout_constraints(shape, text_frame):
    try:
        text_frame.word_wrap = True
        if shape and not shape.has_table:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except: pass

def clean_json_string(raw_text):
    raw_text = raw_text.replace("```json", "").replace("```", "")
    start_idx = raw_text.find('[')
    end_idx = raw_text.rfind(']')
    if start_idx == -1 or end_idx == -1: return raw_text
    return raw_text[start_idx : end_idx + 1]

def get_batch_translation(text_list, target_language, model):
    clean_list = [t for t in text_list if t and len(t.strip()) > 1]
    if not clean_list: return {}
    prompt = f"""
    You are a professional translator. Translate this list to {target_language}.
    Return ONLY a raw JSON list of strings. Do not use Markdown.
    Input: {json.dumps(clean_list, ensure_ascii=False)}
    """
    try:
        response = model.generate_content(prompt)
        cleaned_response = clean_json_string(response.text)
        translated_list = json.loads(cleaned_response)
        if len(clean_list) == len(translated_list):
            return dict(zip(clean_list, translated_list))
    except: pass
    return {}

def extract_image_text(image_blob, target_language, model):
    prompt = f"Extract all text from this image and translate it to {target_language}. Return only the translated text as a single string."
    try:
        response = model.generate_content([prompt, {"mime_type": "image/png", "data": image_blob}])
        return response.text.strip()
    except: return None

# --- MAIN APP ---
if api_key:
    genai.configure(api_key=api_key)
    try:
        model = genai.GenerativeModel('gemini-2.5-flash')
    except:
        model = genai.GenerativeModel('gemini-1.5-flash')

    st.write("### Upload Presentation")
    
    # The Custom CSS above will automatically apply to this file_uploader
    uploaded_file = st.file_uploader("Choose .pptx", type="pptx")
    
    col1, col2, col3 = st.columns(3)
    with col1:
        target_lang = st.selectbox("Target Language", ["English", "Korean", "Japanese", "German", "French", "Chinese"])
    with col2:
        target_font = st.selectbox("Font Family", ["Montserrat", "Arial", "Calibri", "Malgun Gothic"])
    with col3:
        replace_images = st.checkbox("Translate Images?", value=False)

    if uploaded_file and st.button("Translate & Keep Formatting"):
        
        prs = Presentation(uploaded_file)
        
        # --- UI KICKSTART ANIMATION (0% -> 10%) ---
        progress_bar = st.progress(0.0)
        status_text = st.empty()
        status_text.text("Initializing AI Engine...")
        
        # Smooth ignition over 1 second
        for i in range(11):
            time.sleep(0.08) # 0.08 * 10 = 0.8 seconds approx
            progress_bar.progress(i / 100)
            
        # --- CALCULATE WORKLOAD ---
        total_slides = len(prs.slides)
        total_shapes_count = 0
        for s in prs.slides:
            total_shapes_count += len(s.shapes)
            
        # Total "Units" of work (roughly estimated)
        # We start at 10% (0.10)
        # We have 90% (0.90) left to fill
        current_progress = 0.10
        
        # Determine increment per slide
        # Each slide has 3 phases: Harvest (Fast), Translate (Slow), Apply (Fast but granular)
        # We will allocate the 0.90 remaining space divided by total slides
        progress_per_slide = 0.85 / total_slides # Reserve last 5% for cleanup
        
        for i, slide in enumerate(prs.slides):
            
            # Phase 1: Harvest (Tiny creep)
            status_text.text(f"Slide {i+1}/{total_slides}: Analyzing layout...")
            slide_texts = []
            paragraphs_to_process = []
            
            for shape, text_frame in iter_text_frames(slide.shapes):
                enforce_layout_constraints(shape, text_frame)
                for paragraph in text_frame.paragraphs:
                    txt = paragraph.text.strip()
                    if len(txt) > 1:
                        slide_texts.append(txt)
                        paragraphs_to_process.append(paragraph)
            
            # Tiny movement to show we are alive
            current_progress += 0.01
            progress_bar.progress(min(current_progress, 0.95))

            # Phase 2: Translate (The "Fake Upload" Creep)
            if slide_texts:
                status_text.text(f"Slide {i+1}/{total_slides}: Connecting to Gemini...")
                
                # FAKE "Sending Data" Animation before blocking call
                for _ in range(3): 
                    time.sleep(0.1)
                    current_progress += 0.005 # Tiny increments
                    progress_bar.progress(min(current_progress, 0.95))
                    
                translation_map = get_batch_translation(slide_texts, target_lang, model)
                
                # BIG JUMP after blocking call succeeds
                current_progress += (progress_per_slide * 0.4) # 40% of this slide's budget
                progress_bar.progress(min(current_progress, 0.95))

            # Phase 3: Apply (Granular Fill)
            if slide_texts and translation_map:
                status_text.text(f"Slide {i+1}/{total_slides}: Formatting text...")
                
                current_slide_scale = manual_scale
                if scaling_mode == "Smart Adaptive Curve":
                    current_slide_scale = calculate_slide_global_ratio(translation_map)

                items_count = len(paragraphs_to_process)
                if items_count > 0:
                    # Allocate 40% of this slide's budget to formatting
                    format_budget = progress_per_slide * 0.4
                    step_size = format_budget / items_count
                    
                    for idx, paragraph in enumerate(paragraphs_to_process):
                        original = paragraph.text.strip()
                        if original in translation_map:
                            translated_text = translation_map[original]
                            if paragraph.runs:
                                original_run = paragraph.runs[0]
                                paragraph.clear() 
                                new_run = paragraph.add_run() 
                                new_run.text = translated_text
                                apply_scaling(original_run, new_run, target_font, current_slide_scale)
                        
                        # SMOOTH TICK for every paragraph
                        current_progress += step_size
                        progress_bar.progress(min(current_progress, 0.98))

            # Phase 4: Layouts & Images (Remaining budget)
            for shape, text_frame in iter_text_frames(slide.shapes):
                if shape and not shape.has_table:
                      fix_text_overflow_centered(shape, text_frame.text)

            for shape in slide.shapes:
                if shape.has_table:
                    auto_widen_columns(shape)
            
            if replace_images:
                status_text.text(f"Slide {i+1}/{total_slides}: Processing images...")
                images_to_process = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images_to_process.append(shape)
                
                for pic in images_to_process:
                    # Tiny tick for each image
                    current_progress += 0.005
                    progress_bar.progress(min(current_progress, 0.99))
                    try:
                        translated_img_text = extract_image_text(pic.image.blob, target_lang, model)
                        if translated_img_text:
                            left, top = pic.left, pic.top
                            width, height = pic.width, pic.height
                            pic._element.getparent().remove(pic._element)
                            new_tx_box = slide.shapes.add_textbox(left, top, width, height)
                            tf = new_tx_box.text_frame
                            tf.text = translated_img_text
                            for p in tf.paragraphs:
                                p.font.name = target_font
                                p.font.size = Pt(10)
                    except: pass

        # Final Satisfaction
        progress_bar.progress(1.0)
        status_text.success("Translation Complete!")
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        st.download_button("Download Result", out, f"Translated_{uploaded_file.name}")

else:

    st.warning("Please enter API Key in the sidebar.")

